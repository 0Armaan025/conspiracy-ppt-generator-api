from flask import Flask, jsonify, request
import anthropic
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
import logging
import google.generativeai as genai
import replicate
import requests

load_dotenv()

app = Flask(__name__)

client = anthropic.Anthropic(api_key=os.getenv('ANTHROPIC_API_KEY'))
gemini_api_key = os.getenv('GEMINI_API_KEY')
genai.configure(api_key=gemini_api_key)
replicate_api_token = os.getenv('REPLICATE_API_TOKEN')
replicate.Client(api_token=replicate_api_token)

slides_preset = ['title', 'bullets', 'bullets', 'text', 'image', 'thanks']

logging.basicConfig(level=logging.DEBUG)

def call_gemini(prompt):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(prompt)
    text = response.text
    text = text.replace('*', ' ')
    return text

def generate_image(prompt, file_path):
    input = {
        "prompt": prompt,
        "width": 768,
        "height": 768,
        "refine": "expert_ensemble_refiner",
        "apply_watermark": False,
        "num_inference_steps": 25
    }
    output = replicate.run(
        "stability-ai/sdxl:7762fd07cf82c948538e41f63f77d685e02b063e37e496e96eefd46c929f9bdc",
        input=input
    )
    image_url = output[0]

    # Download the image
    if image_url:
        response = requests.get(image_url)
        if response.status_code == 200:
            with open(file_path, 'wb') as f:
                f.write(response.content)

bullets = []
text_content = None
title_slide1 = ""

def get_slide_content(slide_number, theory):
    global bullets
    global text_content
    global title_slide1

    preset = slides_preset[slide_number]

    if preset == 'title':
        title_slide1 = call_gemini(f'Create an amazing title for the presentation about the theory: {theory}')
        intro_paragraph = call_gemini(f'Write a short 40-word introduction about the theory: {theory}. Use easy English, include some real evidence and some false made-up evidence, just for entertainment.')
        return title_slide1, intro_paragraph
    elif preset == 'bullets':
        bullets_title = title_slide1  # Use the same title as the first slide
        num_bullets = 2  # 2 bullets on slide 2 and 3
        bullets = []
        for i in range(num_bullets):
            bullet = call_gemini(f'Create a bullet point to prove this theory right in a fun way of course, anything can happen ;) , theory: {theory}. Use easy English, include some real evidence and some false made-up evidence, just for entertainment, make sure the word limit is 30 for it, something fun that would also make them a bit of laugh and convince them that it is the truth')
            bullets.append(bullet.replace('-', ' '))
        return bullets_title, bullets
    elif preset == 'text':
        text_content = call_gemini(f'Write a paragraph for the paragraph slide on the theory: {theory}. Use easy English, include some real evidence and some false made-up evidence, just for entertainment.')
        return None, text_content
    elif preset == 'image':
        image_prompt = f"Generate an image related to the theory: {theory}"
        image_path = f"image_{slide_number}.png"
        generate_image(image_prompt, image_path)
        return None, image_path
    elif preset == 'thanks':
        return 'Thank You!', None
    return None, None

def get_theory():
    try:
        model = genai.GenerativeModel('gemini-pro')
        response = model.generate_content("You are a conspiracy scientist who writes about conspiracies just for fun, and not in a serious way for kids books. Please suggest a conspiracy theory idea for a movie, only one string. Use easy English and some fun theory ")
        text = response.text
        text = text.replace('*', ' ')
        return text
    except Exception as e:
        logging.error(f"Error getting theory: {e}")
        return "Default Conspiracy Theory"

@app.route('/')
def home():
    theory = get_theory()
    return theory

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    theory = request.form.get('theory')  # Retrieve theory from the POST request

    if not theory:
        theory = get_theory()  # Use AI to generate a default theory if none provided

    prs = Presentation()

    for i, preset in enumerate(slides_preset):
        slide_layout = prs.slide_layouts[1] if preset != 'title' else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title, content = get_slide_content(i, theory)

        if preset == 'title':
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.size = Pt(20)
            
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.text = content
            text_frame.paragraphs[0].font.size = Pt(16)
        elif preset == 'bullets':
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(18)  # Set the title font size for the bullet slide
            text_frame = slide.shapes.placeholders[1].text_frame
            slide.shapes.placeholders[1].top = Inches(1.85)  # Move the bullets placeholder up
            slide.shapes.placeholders[1].width = Inches(9.5)  # Adjust the width of the text box
            for bullet in content:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(15)
        elif preset == 'text':
            slide.shapes.title.text = "Details"
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.text = content
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
        elif preset == 'image':
            slide.shapes.title.text = "Illustration"
            img_path = content  # content contains the path to the generated image
            img = slide.shapes.add_picture(img_path, Inches(2), Inches(2), width=Inches(6), height=Inches(4.5))
        elif preset == 'thanks':
            slide.shapes.title.text = title
            left = Inches(2)
            top = Inches(3)
            width = Inches(6)
            height = Inches(1.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = "Thank You!"
            p.font.size = Pt(32)
            p.font.bold = True
            p.shadow = True

    prs.save('generated_presentation.pptx')

    return jsonify({'message': 'Presentation generated and saved to local storage as generated_presentation.pptx'})


if __name__ == '__main__':
    app.run(debug=True)
